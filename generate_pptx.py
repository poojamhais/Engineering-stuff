"""
Generate PPT and PDF Report for Ohm's Law Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"c:\Users\sanji\Downloads\Presentation"
ASSETS = os.path.join(BASE, "assets")

# Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x1A)
CYAN = RGBColor(0x06, 0xD6, 0xA0)
GOLD = RGBColor(0xFF, 0xD1, 0x66)
PINK = RGBColor(0xEF, 0x47, 0x6F)
BLUE = RGBColor(0x11, 0x8A, 0xB2)
WHITE = RGBColor(0xF0, 0xF4, 0xF8)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0x16, 0x1E, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_after=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_after = space_after
    return p

# ═══════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Top accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = CYAN; shape.line.fill.background()

add_text(slide, "MATOSHRI EDUCATION SOCIETY'S", Inches(0), Inches(0.6), Inches(13.333), Inches(0.5),
         font_size=16, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "MIT POLYTECHNIC & ENGINEERING, YEOLA", Inches(0), Inches(1.0), Inches(13.333), Inches(0.5),
         font_size=20, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Department of Electrical Engineering", Inches(0), Inches(1.45), Inches(13.333), Inches(0.4),
         font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, "OHM'S LAW", Inches(0), Inches(2.2), Inches(13.333), Inches(1.2),
         font_size=72, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Formula box
box = add_shape_bg(slide, Inches(4.5), Inches(3.5), Inches(4.333), Inches(0.9))
add_text(slide, "V = I × R", Inches(4.5), Inches(3.55), Inches(4.333), Inches(0.8),
         font_size=40, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Student cards
names = [("Pooja Santosh Mhais", "Roll No. 1"), ("Rutuja Ramesh Gaikwad", "Roll No. 2"), ("Shubham Madhukar Gaikwad", "Roll No. 3")]
for i, (name, roll) in enumerate(names):
    x = Inches(2.2 + i * 3.2)
    add_shape_bg(slide, x, Inches(4.8), Inches(2.8), Inches(0.9))
    add_text(slide, name, x, Inches(4.85), Inches(2.8), Inches(0.45), font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, roll, x, Inches(5.25), Inches(2.8), Inches(0.35), font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, "Under the Guidance of Miss. Beldar G.N.", Inches(0), Inches(6.0), Inches(13.333), Inches(0.4),
         font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, "Academic Year: 2025-2026", Inches(0), Inches(6.4), Inches(13.333), Inches(0.4),
         font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 2: WHAT IS OHM'S LAW
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "What is Ohm's Law?", Inches(0.8), Inches(0.4), Inches(7), Inches(0.7),
         font_size=36, color=CYAN, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.5), Inches(2.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Ohm's Law is the most fundamental relationship in electrical engineering. Discovered by Georg Simon Ohm in 1827, it describes how voltage, current, and resistance are interrelated in an electrical circuit."
p.font.size = Pt(16); p.font.color.rgb = MUTED; p.font.name = 'Calibri'; p.space_after = Pt(12)
add_para(tf, "V = I × R", font_size=32, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT, space_after=Pt(16))

# V, I, R cards
labels = [("⚡ Voltage (V)", "The electrical pressure that pushes electrons through a conductor. Measured in Volts (V)."),
          ("🌊 Current (I)", "The rate of flow of electric charge through a conductor. Measured in Amperes (A)."),
          ("🔥 Resistance (R)", "The opposition to the flow of current in a circuit. Measured in Ohms (Ω).")]
for i, (title, desc) in enumerate(labels):
    y = Inches(3.6 + i * 1.15)
    add_shape_bg(slide, Inches(0.8), y, Inches(6.5), Inches(1.0))
    add_text(slide, title, Inches(1.0), y + Inches(0.08), Inches(6), Inches(0.35), font_size=16, color=CYAN, bold=True)
    add_text(slide, desc, Inches(1.0), y + Inches(0.45), Inches(6), Inches(0.5), font_size=12, color=MUTED)

# Image
img_path = os.path.join(ASSETS, "georg_ohm.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(8.2), Inches(0.8), Inches(4.5))

# ═══════════════════════════════════════════
# SLIDE 3: BASIC CIRCUIT
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "The Basic Circuit", Inches(0.8), Inches(0.4), Inches(7), Inches(0.7),
         font_size=36, color=CYAN, bold=True)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(6.5), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "A simple circuit consists of a voltage source (battery), a load (resistor), and connecting wires. Current flows from the positive terminal through the circuit and back to the negative terminal."
p.font.size = Pt(15); p.font.color.rgb = MUTED; p.font.name = 'Calibri'

cards = [("🔋 Voltage Source", "Provides the electromotive force (EMF) to drive current through the circuit."),
         ("📐 Resistor", "Opposes current flow. Color bands indicate its resistance value in Ohms.")]
for i, (t, d) in enumerate(cards):
    y = Inches(3.0 + i * 1.15)
    add_shape_bg(slide, Inches(0.8), y, Inches(6.5), Inches(1.0))
    add_text(slide, t, Inches(1.0), y + Inches(0.08), Inches(6), Inches(0.35), font_size=16, color=GOLD, bold=True)
    add_text(slide, d, Inches(1.0), y + Inches(0.45), Inches(6), Inches(0.5), font_size=12, color=MUTED)

add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(6.5), Inches(0.8))
add_text(slide, "Key Insight: If you increase the voltage while keeping resistance constant, the current increases proportionally.", Inches(1.0), Inches(5.35), Inches(6.2), Inches(0.7), font_size=13, color=CYAN)

img_path = os.path.join(ASSETS, "circuit_diagram.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(8.2), Inches(0.8), Inches(4.5))

# ═══════════════════════════════════════════
# SLIDE 4: WATER ANALOGY
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "The Water Analogy", Inches(0.8), Inches(0.4), Inches(7), Inches(0.7),
         font_size=36, color=CYAN, bold=True)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.5), Inches(1.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The easiest way to understand Ohm's Law is by comparing electricity to water flowing through pipes."
p.font.size = Pt(15); p.font.color.rgb = MUTED; p.font.name = 'Calibri'

analogies = [("Voltage = Water Pressure", "Like water pressure pushing water through a pipe, voltage pushes electrons through a wire.", GOLD),
             ("Current = Water Flow Rate", "The amount of water flowing per second is like current — the amount of charge flowing per second.", CYAN),
             ("Resistance = Pipe Width", "A narrow pipe restricts water flow, just as high resistance restricts electric current.", PINK)]
for i, (t, d, c) in enumerate(analogies):
    y = Inches(2.4 + i * 1.3)
    add_shape_bg(slide, Inches(0.8), y, Inches(6.5), Inches(1.1))
    add_text(slide, t, Inches(1.0), y + Inches(0.08), Inches(6), Inches(0.35), font_size=16, color=c, bold=True)
    add_text(slide, d, Inches(1.0), y + Inches(0.48), Inches(6), Inches(0.55), font_size=12, color=MUTED)

img_path = os.path.join(ASSETS, "water_analogy.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(8.2), Inches(0.8), Inches(4.5))

# ═══════════════════════════════════════════
# SLIDE 5: VIR TRIANGLE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "The Ohm's Law Triangle", Inches(0), Inches(0.3), Inches(13.333), Inches(0.7),
         font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Cover the variable you want to find — the remaining arrangement shows the formula!", Inches(0), Inches(1.0), Inches(13.333), Inches(0.5),
         font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

# Triangle shape
shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.5), Inches(1.8), Inches(4.3), Inches(3.5))
shape.fill.background()
shape.line.color.rgb = CYAN; shape.line.width = Pt(2)

add_text(slide, "V", Inches(6.1), Inches(2.5), Inches(1.2), Inches(0.8), font_size=48, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "I", Inches(4.8), Inches(4.2), Inches(1.2), Inches(0.8), font_size=48, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "×", Inches(6.1), Inches(4.2), Inches(1.2), Inches(0.8), font_size=36, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "R", Inches(7.3), Inches(4.2), Inches(1.2), Inches(0.8), font_size=48, color=PINK, bold=True, alignment=PP_ALIGN.CENTER)

# Formulas below
formulas = [("Find Voltage:", "V = I × R", GOLD), ("Find Current:", "I = V / R", CYAN), ("Find Resistance:", "R = V / I", PINK)]
for i, (label, formula, c) in enumerate(formulas):
    x = Inches(2 + i * 3.5)
    add_shape_bg(slide, x, Inches(5.7), Inches(3), Inches(0.9))
    add_text(slide, label, x, Inches(5.72), Inches(3), Inches(0.35), font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_text(slide, formula, x, Inches(6.05), Inches(3), Inches(0.5), font_size=22, color=c, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 6: INTERACTIVE SIMULATOR (Static version)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "Ohm's Law Simulator", Inches(0), Inches(0.3), Inches(13.333), Inches(0.7),
         font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Adjust voltage and resistance to see how current changes (Interactive version available in HTML)", Inches(0), Inches(1.0), Inches(13.333), Inches(0.5),
         font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# Example calculations
examples = [("12V / 100Ω", "= 0.120 A", "1.44 W"), ("24V / 8Ω", "= 3.000 A", "72.00 W"), ("120V / 60Ω", "= 2.000 A", "240.00 W")]
for i, (inp, curr, pwr) in enumerate(examples):
    x = Inches(1 + i * 4)
    add_shape_bg(slide, x, Inches(1.8), Inches(3.5), Inches(4.5))
    add_text(slide, f"Example {i+1}", x, Inches(1.9), Inches(3.5), Inches(0.4), font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, inp, x, Inches(2.4), Inches(3.5), Inches(0.5), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, "Current", x, Inches(3.2), Inches(3.5), Inches(0.3), font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_text(slide, curr, x, Inches(3.5), Inches(3.5), Inches(0.7), font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, "Power", x, Inches(4.5), Inches(3.5), Inches(0.3), font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_text(slide, pwr, x, Inches(4.8), Inches(3.5), Inches(0.6), font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 7: SOLVED EXAMPLES
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "Solved Examples", Inches(0), Inches(0.3), Inches(13.333), Inches(0.7),
         font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

problems = [
    ("Example 1: Finding Current", "Given: V = 24V, R = 8Ω", "I = V / R = 24 / 8", "I = 3 Amperes"),
    ("Example 2: Finding Voltage", "Given: I = 5A, R = 12Ω", "V = I × R = 5 × 12", "V = 60 Volts"),
    ("Example 3: Finding Resistance", "Given: V = 120V, I = 0.5A", "R = V / I = 120 / 0.5", "R = 240 Ohms"),
]
for i, (title, given, solution, answer) in enumerate(problems):
    y = Inches(1.2 + i * 2.0)
    add_shape_bg(slide, Inches(1.5), y, Inches(10.333), Inches(1.75))
    add_text(slide, title, Inches(1.8), y + Inches(0.08), Inches(9.5), Inches(0.35), font_size=18, color=GOLD, bold=True)
    add_text(slide, given, Inches(1.8), y + Inches(0.45), Inches(9.5), Inches(0.3), font_size=14, color=CYAN, bold=True)
    add_text(slide, f"Solution: {solution}", Inches(1.8), y + Inches(0.8), Inches(9.5), Inches(0.3), font_size=14, color=MUTED)
    add_text(slide, f"Answer: {answer}", Inches(1.8), y + Inches(1.2), Inches(9.5), Inches(0.35), font_size=16, color=CYAN, bold=True)

# ═══════════════════════════════════════════
# SLIDE 8: APPLICATIONS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "Practical Applications", Inches(0.8), Inches(0.4), Inches(7), Inches(0.7),
         font_size=36, color=CYAN, bold=True)

apps = [
    ("💡 LED Circuit Design", "Calculate the correct resistor value to safely drive an LED without burning it out."),
    ("🔌 Power Supply Design", "Determine current draw and power dissipation in voltage regulator circuits."),
    ("🏠 Household Wiring", "Calculate wire gauge requirements and fuse ratings for safe home electrical installations."),
    ("🔧 Troubleshooting", "Diagnose faulty components by measuring voltage drops and comparing to expected values."),
]
for i, (t, d) in enumerate(apps):
    y = Inches(1.3 + i * 1.3)
    add_shape_bg(slide, Inches(0.8), y, Inches(6.5), Inches(1.1))
    add_text(slide, t, Inches(1.0), y + Inches(0.08), Inches(6), Inches(0.35), font_size=16, color=GOLD, bold=True)
    add_text(slide, d, Inches(1.0), y + Inches(0.48), Inches(6), Inches(0.5), font_size=12, color=MUTED)

img_path = os.path.join(ASSETS, "applications.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(8.2), Inches(0.8), Inches(4.5))

# ═══════════════════════════════════════════
# SLIDE 9: LIMITATIONS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, "Limitations of Ohm's Law", Inches(0), Inches(0.3), Inches(13.333), Inches(0.7),
         font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Ohm's Law applies to ohmic (linear) conductors under constant physical conditions.", Inches(0), Inches(1.0), Inches(13.333), Inches(0.5),
         font_size=15, color=MUTED, alignment=PP_ALIGN.CENTER)

limits = [
    "Non-linear devices — Diodes, transistors, and LEDs do not obey Ohm's Law.",
    "Temperature dependence — Resistance changes with temperature.",
    "High-frequency AC circuits — Impedance must be considered.",
    "Semiconductors — Complex V-I characteristics vary with doping.",
    "Superconductors — Zero resistance at extremely low temperatures.",
]
for i, lim in enumerate(limits):
    y = Inches(1.8 + i * 1.05)
    add_shape_bg(slide, Inches(1.5), y, Inches(10.333), Inches(0.85))
    add_text(slide, f"⚠️  {lim}", Inches(1.8), y + Inches(0.15), Inches(9.8), Inches(0.55), font_size=15, color=MUTED)

# ═══════════════════════════════════════════
# SLIDE 10: THANK YOU
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, "Thank You!", Inches(0), Inches(1.5), Inches(13.333), Inches(1.5),
         font_size=72, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "We hope this presentation helped you understand the fundamentals of Ohm's Law.", Inches(0), Inches(3.2), Inches(13.333), Inches(0.6),
         font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

box = add_shape_bg(slide, Inches(4.5), Inches(4.0), Inches(4.333), Inches(0.9))
add_text(slide, "V = I × R", Inches(4.5), Inches(4.05), Inches(4.333), Inches(0.8),
         font_size=40, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

for i, (name, roll) in enumerate(names):
    x = Inches(2.2 + i * 3.2)
    add_shape_bg(slide, x, Inches(5.2), Inches(2.8), Inches(0.9))
    add_text(slide, name, x, Inches(5.25), Inches(2.8), Inches(0.45), font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, roll, x, Inches(5.65), Inches(2.8), Inches(0.35), font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, "MIT Polytechnic & Engineering, Yeola • Electrical Engineering", Inches(0), Inches(6.4), Inches(13.333), Inches(0.4),
         font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = CYAN; shape.line.fill.background()

# SAVE
pptx_path = os.path.join(BASE, "Ohms_Law_Presentation.pptx")
prs.save(pptx_path)
print(f"PPTX saved: {pptx_path}")
